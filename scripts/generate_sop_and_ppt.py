"""
Drug Management System (DMS) — SOP & Adobe Sign PPT Generator
Produces:
  release/assets/DMS_SOP_OperationalGuide.docx
  release/assets/DMS_AdobeSign_Guide.pptx

Run: python3 scripts/generate_sop_and_ppt.py
"""

import os
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as Inch
import datetime

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "release", "assets")
os.makedirs(OUTPUT_DIR, exist_ok=True)
DOCX_PATH = os.path.join(OUTPUT_DIR, "DMS_SOP_OperationalGuide.docx")
PPTX_PATH = os.path.join(OUTPUT_DIR, "DMS_AdobeSign_Guide.pptx")

TODAY = datetime.date.today().strftime("%B %d, %Y")
VERSION = "1.0"

# ─────────────────────────────────────────────────────────────────────────────
# COLOUR PALETTE
# ─────────────────────────────────────────────────────────────────────────────
BLUE     = RGBColor(0x1A, 0x23, 0x7E)   # dark blue headings
ACCENT   = RGBColor(0x19, 0x76, 0xD2)   # mid blue
TEAL     = RGBColor(0x00, 0x79, 0x6B)   # teal for tips
GRAY     = RGBColor(0x61, 0x61, 0x61)   # body text
LIGHT    = RGBColor(0xF5, 0xF5, 0xF5)   # row shading


# ─────────────────────────────────────────────────────────────────────────────
# DOCX HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _shade_table_header(row):
    """Apply blue background to a table header row."""
    for cell in row.cells:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), "1A237E")
        tcPr.append(shd)
        for para in cell.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.bold = True


def _set_cell_bg(cell, hex_color: str):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def add_heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    run = p.runs[0] if p.runs else p.add_run(text)
    run.font.color.rgb = BLUE if level == 1 else ACCENT
    run.font.bold = True
    p.paragraph_format.space_before = Pt(14 if level == 1 else 8)
    p.paragraph_format.space_after = Pt(4)
    return p


def add_body(doc, text):
    p = doc.add_paragraph(text)
    p.paragraph_format.space_after = Pt(4)
    for run in p.runs:
        run.font.size = Pt(10.5)
        run.font.color.rgb = GRAY
    return p


def add_step(doc, number, text):
    p = doc.add_paragraph(style="List Number")
    run = p.add_run(f"Step {number}: ")
    run.bold = True
    run.font.color.rgb = ACCENT
    run2 = p.add_run(text)
    run2.font.size = Pt(10.5)
    run2.font.color.rgb = GRAY
    p.paragraph_format.space_after = Pt(3)
    return p


def add_bullet(doc, text):
    p = doc.add_paragraph(style="List Bullet")
    run = p.add_run(text)
    run.font.size = Pt(10.5)
    run.font.color.rgb = GRAY
    p.paragraph_format.space_after = Pt(2)
    return p


def add_note(doc, text):
    p = doc.add_paragraph()
    run = p.add_run("Note: ")
    run.bold = True
    run.font.color.rgb = TEAL
    run2 = p.add_run(text)
    run2.font.size = Pt(10)
    run2.font.color.rgb = TEAL
    p.paragraph_format.left_indent = Pt(18)
    p.paragraph_format.space_after = Pt(4)
    return p


def add_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    hdr = table.rows[0]
    for i, h in enumerate(headers):
        hdr.cells[i].text = h
        hdr.cells[i].paragraphs[0].runs[0].font.size = Pt(10)
    _shade_table_header(hdr)
    for r_idx, row in enumerate(rows):
        tr = table.rows[r_idx + 1]
        bg = "F5F5F5" if r_idx % 2 == 0 else "FFFFFF"
        for c_idx, val in enumerate(row):
            cell = tr.cells[c_idx]
            cell.text = val
            cell.paragraphs[0].runs[0].font.size = Pt(10)
            _set_cell_bg(cell, bg)
    doc.add_paragraph()
    return table


def add_separator(doc):
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "4")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), "1A237E")
    pBdr.append(bottom)
    pPr.append(pBdr)


# ─────────────────────────────────────────────────────────────────────────────
# WORD DOCUMENT
# ─────────────────────────────────────────────────────────────────────────────

def build_docx():
    doc = Document()

    # ── Default styles ──────────────────────────────────────────────────────
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(10.5)

    # ── Margins ─────────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.2)
        section.right_margin = Inches(1.2)

    # ════════════════════════════════════════════════════════════════════════
    # COVER PAGE
    # ════════════════════════════════════════════════════════════════════════
    doc.add_paragraph("\n\n\n")
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Drug Management System")
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BLUE

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = sub.add_run("User Operational Guide")
    r2.font.size = Pt(18)
    r2.font.color.rgb = ACCENT

    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"Version {VERSION}  |  {TODAY}  |  SharePoint Online")

    doc.add_paragraph("\n")
    audience = doc.add_paragraph()
    audience.alignment = WD_ALIGN_PARAGRAPH.CENTER
    audience.add_run("Audience: Administrator  |  HR  |  Author  |  Approver")

    doc.add_page_break()

    # ════════════════════════════════════════════════════════════════════════
    # 1. INTRODUCTION
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "1. Introduction", 1)
    add_body(doc,
        "The Drug Management System (DMS) is a SharePoint Online application that enables "
        "pharmaceutical teams to create, review, approve, and digitally sign drug-related documents "
        "following a structured, role-based workflow. This guide describes every screen and action "
        "in plain operational terms — telling you exactly what to click, what fields appear, and "
        "what happens next."
    )

    add_heading(doc, "1.1 Role Permissions Overview", 2)
    add_body(doc,
        "Access to each screen is controlled by SharePoint Group membership. "
        "The table below summarises which tabs and screens each role can see and what they can do."
    )
    add_table(doc,
        ["Role / Group", "Screens Accessible", "Permitted Actions"],
        [
            ["Admin",
             "Category, Template, Drug, User, eCTD/CTD Folders, All Documents, Approver Dashboard",
             "Full CRUD on all screens; view all documents; manage users"],
            ["HR",
             "Category, Template, Drug, Documents",
             "Add/edit Categories, Templates, Drugs; view documents (User screen is Admin-only)"],
            ["Author",
             "My Documents, Add Document",
             "Create and submit documents; view own documents and version history"],
            ["Approver",
             "Assigned To Me, Approver Dashboard",
             "Review, approve, or reject documents assigned to them; initiate Adobe Sign"],
        ]
    )
    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 2. CATEGORY SCREEN
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "2. Category Screen (Admin / HR only)", 1)
    add_body(doc,
        "The Category screen provides a hierarchical folder view that organises document categories "
        "across five levels: Document Category → Group → SubGroup → Artifact Name → Template Name. "
        "Each level drills down into the next when you click a row."
    )

    add_heading(doc, "2.1 Navigating the Hierarchy", 2)
    add_body(doc,
        "When the screen loads, you see all top-level Document Categories listed as rows. "
        "A breadcrumb bar at the top shows your current location (e.g., Categories > Regulatory > CTD)."
    )
    add_step(doc, 1, "Click any row to drill into its child level (Group, SubGroup, etc.).")
    add_step(doc, 2, "Click a breadcrumb segment to jump back to that level.")
    add_step(doc, 3, "Use the Search bar to filter items at the current level by name.")
    add_step(doc, 4, "Use the Status filter dropdown (Active / Inactive) to narrow the list.")

    add_heading(doc, "2.2 Adding a Category", 2)
    add_step(doc, 1, "Click the 'Add Category' button at the top-right of the grid.")
    add_step(doc, 2, "The 'Create New Category' form opens. Fill in the following fields:")
    add_table(doc,
        ["Field", "Required?", "Notes"],
        [
            ["Category Name", "Yes (*)", "Free-text label for this category entry."],
            ["Document Category", "Yes (*)", "Dropdown — the top-level document category."],
            ["Group", "No", "Dropdown — the Group within the Document Category."],
            ["SubGroup", "No", "Dropdown — enabled only after Group is selected."],
            ["Artifact Name", "No", "Dropdown — the artifact type."],
            ["Template Name", "No", "Dropdown — the template associated with this artifact."],
            ["Status", "Yes (*)", "Active or Inactive. Defaults to Active."],
        ]
    )
    add_step(doc, 3, "Click 'Add Category' to save. A success banner appears when saved.")
    add_note(doc, "Fields marked * are required. The form will not save if they are missing.")

    add_heading(doc, "2.3 Editing a Category", 2)
    add_step(doc, 1, "Select the row you want to edit (checkbox click), then click the pencil (Edit) icon.")
    add_step(doc, 2, "The same form opens, pre-filled with the existing values.")
    add_step(doc, 3, "Make your changes and click 'Save Changes'.")

    add_heading(doc, "2.4 Deleting a Category", 2)
    add_step(doc, 1, "Select one or more rows using the checkboxes.")
    add_step(doc, 2, "Click the delete (trash) icon in the toolbar.")
    add_step(doc, 3, "A confirmation dialog appears. Click 'Delete' to proceed or 'Cancel' to abort.")
    add_note(doc, "Deleting a category row removes that specific entry. It does not cascade-delete child levels unless those rows are also selected.")

    add_heading(doc, "2.5 Quick-Add Inline Feature", 2)
    add_body(doc,
        "At any hierarchy level, a quick-add button lets you add a new value directly without leaving "
        "the list. Click the '+' icon next to the level header, type the new value, then press Save. "
        "This is faster than opening the full form when you only need to add a single node."
    )

    add_heading(doc, "2.6 Excel Bulk Upload", 2)
    add_step(doc, 1, "Click the Excel icon button near the top of the screen.")
    add_step(doc, 2, "Download the provided template file and fill in your category data row by row.")
    add_step(doc, 3, "Upload the completed file. The system will validate and import the rows.")
    add_note(doc, "Any row with a missing required field will be skipped and listed in the error report.")

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 3. TEMPLATE SCREEN
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "3. Template Screen (Admin / HR only)", 1)
    add_body(doc,
        "Templates are master document files (Word, PDF, Excel, etc.) that Authors use as the basis "
        "for new documents. Each template is mapped to a regulatory filing structure through the "
        "Mapping Type setting."
    )
    add_body(doc,
        "The template list can be filtered using three tab-style filter buttons at the top of the "
        "grid: All, Active, and Inactive. Selecting 'Active' shows only templates with Active status; "
        "'Inactive' shows archived templates; 'All' shows every template regardless of status."
    )
    add_body(doc,
        "Summary cards at the top show: Total Templates, Active Templates, Inactive Templates, and eCTD Mapped count."
    )

    add_heading(doc, "3.1 Uploading a New Template", 2)
    add_step(doc, 1, "Click 'Add Template' in the grid toolbar.")
    add_step(doc, 2, "The 'Upload Template' page opens. Fill in the following fields:")
    add_table(doc,
        ["Field", "Required?", "Notes"],
        [
            ["Template Name", "Yes (*)", "A descriptive name (e.g., Clinical Trial Protocol v3.0)."],
            ["Version", "No", "Defaults to 1.0."],
            ["Category", "No", "Dropdown — links the template to a category."],
            ["Country", "No", "Dropdown — country-specific applicability."],
            ["Status", "No", "Active or Inactive. Defaults to Active."],
            ["Mapping Type", "No", "None, eCTD, GMP, or TMF (see below)."],
            ["Upload File", "Yes (*)", "Drag and drop or click Browse. Accepted: DOC, DOCX, PDF, XLS, XLSX."],
        ]
    )

    add_heading(doc, "3.2 Mapping Type — Conditional Fields", 2)
    add_body(doc,
        "When you select a Mapping Type, additional fields appear automatically:"
    )
    add_table(doc,
        ["Mapping Type", "Extra Fields That Appear"],
        [
            ["None", "No extra fields."],
            ["eCTD", "eCTD Module (Module 1–5 dropdown), CTD Folder (filtered by module), eCTD Section (filtered by module). All three are required."],
            ["GMP", "GMP Model dropdown (required)."],
            ["TMF", "TMF Folder dropdown (required)."],
        ]
    )
    add_note(doc, "For eCTD: select the Module first — the CTD Folder and eCTD Section lists will then filter to show only entries for that module.")

    add_step(doc, 3, "Click 'Save Template'. A success message appears and you are returned to the list.")

    add_heading(doc, "3.3 Editing a Template", 2)
    add_step(doc, 1, "Select the template row, then click the pencil (Edit) icon.")
    add_step(doc, 2, "The same form opens pre-filled. Template Name is read-only in edit mode.")
    add_step(doc, 3, "To replace the file: click 'Remove' next to the current file name, then upload a new one.")
    add_step(doc, 4, "Click 'Update Template' to save.")

    add_heading(doc, "3.4 Deleting a Template", 2)
    add_step(doc, 1, "Select one or more rows, then click the trash icon.")
    add_step(doc, 2, "Confirm in the dialog that appears.")

    add_heading(doc, "3.5 Previewing / Downloading a Template", 2)
    add_step(doc, 1, "Click the eye icon in the Actions column to open a side preview panel.")
    add_step(doc, 2, "Click the download icon to download the file to your computer.")

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 4. DRUG SCREEN
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "4. Drug Screen (Admin / HR only)", 1)
    add_body(doc,
        "The Drug screen stores records of all drugs managed within the system. Each drug is "
        "referenced when an Author creates a document, and drives the Country and Template filter cascade."
    )
    add_body(doc,
        "Summary cards show: Total Drugs, Active, In Development, and Inactive counts."
    )

    add_heading(doc, "4.1 Adding a Drug", 2)
    add_step(doc, 1, "Click 'Add Drug' in the toolbar.")
    add_step(doc, 2, "The Add Drug form opens. Fill in:")
    add_table(doc,
        ["Field", "Required?", "Notes"],
        [
            ["Drug Name", "Yes (*)", "The official drug name."],
            ["Category", "No", "Free-text category label."],
            ["Status", "No", "In Development / Active / Discontinued / Approved."],
            ["Description", "No", "Optional multi-line description."],
        ]
    )
    add_step(doc, 3, "Click 'Add Drug'. The new drug appears in the grid.")

    add_heading(doc, "4.2 Editing a Drug", 2)
    add_step(doc, 1, "Click the pencil icon on the drug row or select the row and click Edit in the toolbar.")
    add_step(doc, 2, "Update the fields and click 'Update Drug'.")

    add_heading(doc, "4.3 Deleting a Drug", 2)
    add_step(doc, 1, "Click the trash icon on the row or select one or more rows and click Delete.")
    add_step(doc, 2, "Confirm the deletion in the dialog.")

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 5. USER SCREEN
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "5. User Screen (Admin only)", 1)
    add_body(doc,
        "The User screen controls which individuals belong to each SharePoint Group. "
        "Adding a user here simultaneously adds them to the corresponding SharePoint Group, "
        "which is the mechanism that controls their access rights across the entire site."
    )
    add_body(doc,
        "Summary cards show: Admins, HR members, and Authors."
    )

    add_heading(doc, "5.1 Adding a User", 2)
    add_step(doc, 1, "Click 'Add User to Group'.")
    add_step(doc, 2, "Fill in the User form:")
    add_table(doc,
        ["Field", "Required?", "Notes"],
        [
            ["Full Name", "Yes (*)", "The user's display name."],
            ["Email Address", "Yes (*)", "Must match their Microsoft 365 account."],
            ["SharePoint Group (Role)", "Yes (*)", "DMS Admins / DMS HR / DMS Authors / DMS Approvers."],
            ["Status", "No", "Active or Inactive."],
            ["Project / Drug Assignment", "No", "Optional — note which project they work on."],
        ]
    )
    add_note(doc, "A helper note appears below the Role field: 'User will be added to this SharePoint Group'.")
    add_step(doc, 3, "Click 'Add to Group'. The user now has the permissions associated with that group.")

    add_heading(doc, "5.2 Editing a User", 2)
    add_step(doc, 1, "Select the user row, click the pencil icon.")
    add_step(doc, 2, "You can update their Role or Status. Click 'Save Changes'.")

    add_heading(doc, "5.3 Deleting / Removing a User", 2)
    add_step(doc, 1, "Select the user row, click the trash icon.")
    add_step(doc, 2, "A confirmation dialog explains that removing this user will revoke their SharePoint Group permissions.")
    add_step(doc, 3, "Click 'Remove' to confirm.")
    add_note(doc, "Removing a user does not delete their Microsoft 365 account — it only removes them from the DMS SharePoint Group.")

    add_heading(doc, "5.4 Searching and Filtering", 2)
    add_bullet(doc, "Use the Search bar to search by name or email.")
    add_bullet(doc, "Use the Role filter dropdown to show only Admins, HR, or Authors.")

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 6. eCTD / CTD FOLDER SCREEN
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "6. eCTD / CTD Folder Screen (Admin only)", 1)
    add_body(doc,
        "The CTD Folder screen manages the regulatory folder hierarchy used to classify and "
        "locate documents within the system. Three mapping structures are supported:"
    )
    add_table(doc,
        ["Mapping Type", "Full Name", "Purpose"],
        [
            ["eCTD", "Electronic Common Technical Document", "Regulatory submission structure across Modules 1–5."],
            ["GMP", "Good Manufacturing Practice", "Quality and manufacturing document classification."],
            ["TMF", "Trial Master File", "Clinical trial document filing structure."],
        ]
    )

    add_heading(doc, "6.1 Navigating the Folder Hierarchy", 2)
    add_body(doc,
        "The screen opens at the root level, showing top-level modules. "
        "Click a folder name to drill into its children. A breadcrumb trail at the top tracks your location."
    )

    add_heading(doc, "6.2 Adding a Folder", 2)
    add_step(doc, 1, "Click 'Add Root Module' (at the root) or 'Add Subfolder' (inside a folder).")
    add_step(doc, 2, "Fill in the Create Folder form:")
    add_table(doc,
        ["Field", "Required?", "Notes"],
        [
            ["Parent Folder", "No", "Automatically pre-filled if you clicked 'Add Subfolder'. Clear it to make a root entry."],
            ["Folder Code", "Yes (*)", "Standard CTD numbering (e.g., 1.1, 2.3.1). Auto-prefixed with parent code when adding a subfolder."],
            ["Folder Name", "Yes (*)", "Descriptive name (e.g., Administrative Information)."],
            ["Description", "No", "Free-text description of what belongs in this folder."],
            ["Sort Order", "No", "Numeric value controlling display order."],
        ]
    )
    add_step(doc, 3, "Click 'Create Folder' to save.")

    add_heading(doc, "6.3 Editing and Deleting Folders", 2)
    add_step(doc, 1, "Click the pencil icon on any folder row to edit it.")
    add_step(doc, 2, "Click the trash icon to delete it. A dialog warns that all subfolders will also be deleted.")
    add_note(doc, "CTD Folders are referenced when creating Templates and when placing Documents within the eCTD structure.")

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 7. DOCUMENT SCREEN — TAB OVERVIEW
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "7. Document Screen — Tab Overview", 1)
    add_body(doc,
        "The Document screen shows different tabs depending on your role:"
    )
    add_table(doc,
        ["Tab Name", "Who Sees It", "What It Shows"],
        [
            ["All Documents", "Admin only", "Every document in the system, regardless of status or author."],
            ["My Documents", "Authors", "Documents created by the currently logged-in user."],
            ["Assigned To Me", "Approvers", "Documents assigned to the logged-in approver with status Pending Approval or In Review."],
        ]
    )

    add_heading(doc, "7.1 CTD Structure Filter (Left Panel)", 2)
    add_body(doc,
        "On the left side of the Document screen, a folder panel lets you browse documents by their "
        "regulatory filing structure. Toggle between eCTD, GMP, and TMF views. Clicking a folder in "
        "the panel filters the document grid to show only documents mapped to that folder."
    )
    add_bullet(doc, "eCTD view: Shows the CTD module and section hierarchy (Module 1 through 5).")
    add_bullet(doc, "GMP view: Shows the GMP model folders.")
    add_bullet(doc, "TMF view: Shows the Trial Master File folder structure.")
    add_bullet(doc, "Click any folder to filter the main document grid to that folder's documents.")

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 8. ADD DOCUMENT
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "8. Add Document (Author / Admin)", 1)
    add_body(doc,
        "Authors and Admins create new documents from the 'My Documents' tab or the 'All Documents' tab."
    )

    add_heading(doc, "8.1 Creating a Document", 2)
    add_step(doc, 1, "Click 'Add Document' in the grid toolbar.")
    add_step(doc, 2,
        "The Create Document page opens. Section 1 'Document Details' contains four cascading fields:"
    )
    add_table(doc,
        ["Field", "Required?", "Notes"],
        [
            ["Drug", "Yes (*)", "Select the drug this document is for. This filters the Country list."],
            ["Country", "Yes (*)", "Enabled after Drug is selected. Filters the Template list."],
            ["Template", "Yes (*)", "Enabled after Country is selected. Lists active templates for that drug/country combination."],
            ["Approver", "Yes (*)", "A dropdown of users in the DMS Approvers group."],
        ]
    )
    add_body(doc,
        "When a Template is selected, the system automatically shows a summary card with:"
    )
    add_bullet(doc, "Document Name — derived from the template's artifact name.")
    add_bullet(doc, "CTD / eCTD Placement — derived from the template's mapping (e.g., 'CTD: 2.3.1').")
    add_bullet(doc, "Category — the category linked to the template.")

    add_step(doc, 3,
        "Section 2 'Initial Comments' — add optional notes. Click '+ Add Comment' to add more fields. "
        "Click the trash icon beside a comment to remove it."
    )
    add_step(doc, 4, "Click 'Save / Submit'. The document is created with status 'Draft'.")
    add_note(doc,
        "The form is not saveable until all four required fields (Drug, Country, Template, Approver) are filled."
    )

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 9. DOCUMENT STATUS WORKFLOW
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "9. Document Status Workflow", 1)
    add_body(doc,
        "This section walks you through the full lifecycle of a document from creation to final signature."
    )
    add_table(doc,
        ["Stage", "Action Taken By", "Resulting Status", "Who Sees It"],
        [
            ["1. Create", "Author creates document", "Draft", "Author (My Documents tab)"],
            ["2. Submit", "Author clicks 'Submit for Review'", "In Review", "Approver (Assigned To Me tab)"],
            ["3a. Approve", "Approver clicks 'Approve'", "Pending Approval", "Both; Adobe eSign is initiated"],
            ["3b. Reject", "Approver clicks 'Reject' (with reason)", "Rejected", "Author sees it back in My Documents"],
            ["4. Sign", "Approver signs in Adobe Acrobat Sign (email/inline)", "Pending for Signature → Signed", "Both"],
            ["5. Final Approve", "Approver clicks 'Final Approve' after signing", "Signed / Final / Approved", "Admin (Signed Documents library)"],
            ["6. Revise (if rejected)", "Author edits and resubmits", "In Review (again)", "Approver (Assigned To Me tab)"],
        ]
    )

    add_heading(doc, "9.1 Detailed Step-by-Step", 2)
    add_step(doc, 1,
        "[Author] Open the document from 'My Documents'. Click 'Submit for Review'. "
        "A confirmation dialog asks you to confirm. Click 'Submit'. Status changes to 'In Review'."
    )
    add_step(doc, 2,
        "[Approver] Log in and go to 'Assigned To Me'. The document appears with status 'In Review'. "
        "Click the document name to open the detail panel."
    )
    add_step(doc, 3,
        "[Approver] Review the document in the inline viewer. Add Reviewer Comments if needed. "
        "Choose one of: 'Approve' (moves to Pending Approval and initiates Adobe eSign), "
        "or 'Reject' (returns to author with a reason)."
    )
    add_step(doc, 4,
        "[Approver — if Approved] Status changes to 'Pending Approval'. "
        "The system initiates an Adobe Acrobat Sign request automatically. "
        "The approver receives an Adobe Sign notification email with a direct signing link."
    )
    add_step(doc, 5,
        "[Approver] Open the Adobe Sign link in the email (or inline panel). "
        "Review the document and place your electronic signature. Click 'Click to Sign'. "
        "Status moves to 'Pending for Signature' / 'Signed'."
    )
    add_step(doc, 6,
        "[Approver] Once the Adobe Sign signature is confirmed, click 'Final Approve' in the DMS "
        "to mark the document as fully approved. Status updates to 'Signed' / 'Final' / 'Approved'. "
        "The signed PDF is saved to the SharePoint Signed Documents library."
    )
    add_note(doc,
        "At every stage, the document status badge in the grid updates in real time. "
        "The correct tab (My Documents, Assigned To Me, All Documents) will reflect the new status automatically."
    )

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 10. APPROVER DASHBOARD
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "10. Approver Dashboard", 1)
    add_body(doc,
        "When an Approver logs in, they see the Approval Queue (also accessible from the Approver Dashboard menu item). "
        "This view is dedicated to documents that need their attention."
    )

    add_heading(doc, "10.1 What You See", 2)
    add_body(doc, "Summary cards at the top of the page show:")
    add_bullet(doc, "Total Pending — all documents assigned to you awaiting action.")
    add_bullet(doc, "Pending Approval — documents at 'Pending Approval' stage.")
    add_bullet(doc, "In Review — documents you have not yet actioned since they were submitted.")

    add_heading(doc, "10.2 The Document Grid", 2)
    add_body(doc, "The grid lists documents with the following columns:")
    add_table(doc,
        ["Column", "Description"],
        [
            ["Document Name", "The document title, derived from its template artifact name."],
            ["Drug", "The drug this document relates to."],
            ["Category", "The document category."],
            ["Submitted By", "The Author who created and submitted the document."],
            ["Status", "Current workflow status (In Review or Pending Approval)."],
            ["Version", "The version number (e.g., v1)."],
            ["Date", "The date the document was last modified."],
        ]
    )

    add_heading(doc, "10.3 Reviewing a Document", 2)
    add_step(doc, 1, "Click on a document row to open the document detail panel or page.")
    add_step(doc, 2, "The document file is displayed in an inline viewer.")
    add_step(doc, 3, "Add Reviewer Comments in the comments section if needed.")
    add_step(doc, 4, "Click 'Approve' to move the document to 'Pending Approval' status — this initiates the Adobe Acrobat Sign request automatically.")
    add_step(doc, 5, "Complete the signature in the Adobe Sign email or inline panel, then click 'Final Approve' in DMS to set the document to its final 'Signed / Final / Approved' state.")
    add_step(doc, 6, "Alternatively, click 'Reject' to return the document to the Author. You must provide a rejection reason.")
    add_note(doc,
        "Use the Status filter dropdown above the grid to show only 'Pending Approval' or 'In Review' documents."
    )

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # 11. AUTHOR DASHBOARD
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "11. Author Dashboard (My Documents)", 1)
    add_body(doc,
        "The Author Dashboard shows all documents created by the currently logged-in Author."
    )

    add_heading(doc, "11.1 Summary Cards", 2)
    add_bullet(doc, "Total Documents — all documents belonging to you.")
    add_bullet(doc, "Drafts — documents not yet submitted.")
    add_bullet(doc, "Pending Approval — documents you have submitted that are awaiting approval.")
    add_bullet(doc, "Approved — documents that have been approved and/or signed.")
    add_bullet(doc, "Rejected — documents returned to you for revision (shown only if count > 0).")

    add_heading(doc, "11.2 The Document Grid", 2)
    add_body(doc,
        "The grid shows: Document Name, Drug, Category, Status, Version, and Last Modified. "
        "Use the Status filter to narrow the view. Use the Search bar to search by name, drug, or category."
    )

    add_heading(doc, "11.3 Viewing Document Details and History", 2)
    add_step(doc, 1, "Click a document name to open its detail panel.")
    add_step(doc, 2, "In the detail panel you can: view the document file inline, read Reviewer Comments, and see Version History.")
    add_step(doc, 3,
        "If a document is in 'Draft' or 'Rejected' status, the inline viewer allows editing. "
        "If it is 'In Review' or later, the viewer is read-only."
    )

    add_separator(doc)

    # ════════════════════════════════════════════════════════════════════════
    # GLOSSARY
    # ════════════════════════════════════════════════════════════════════════
    add_heading(doc, "Glossary", 1)
    add_table(doc,
        ["Term", "Definition"],
        [
            ["eCTD", "Electronic Common Technical Document — the globally recognised format for regulatory submissions, organised across Modules 1–5."],
            ["GMP", "Good Manufacturing Practice — regulatory standard for the pharmaceutical manufacturing environment."],
            ["TMF", "Trial Master File — the collection of documents that allows a clinical trial to be reconstructed and evaluated."],
            ["Adobe Acrobat Sign", "The electronic signature service used by DMS to obtain legally compliant digital signatures on approved documents."],
            ["Draft", "Initial status of every newly created document."],
            ["In Review", "Status after the Author submits the document; the Approver can now act on it."],
            ["Pending Approval", "Status after the Approver clicks 'Approve'; Adobe Acrobat Sign is initiated."],
            ["Pending for Signature", "Intermediate status while the document is with Adobe Sign awaiting the electronic signature."],
            ["Rejected", "Status after the Approver rejects a document; returned to the Author for revision."],
            ["Signed / Final", "Status after the Adobe Sign signature is completed; document is locked."],
        ]
    )

    doc.save(DOCX_PATH)
    print(f"[OK] SOP Word document saved to: {DOCX_PATH}")


# ─────────────────────────────────────────────────────────────────────────────
# POWERPOINT HELPERS
# ─────────────────────────────────────────────────────────────────────────────

PPTX_W = Inches(13.33)
PPTX_H = Inches(7.5)

PPTX_DARK  = PPTXColor(0x1A, 0x23, 0x7E)
PPTX_BLUE  = PPTXColor(0x19, 0x76, 0xD2)
PPTX_TEAL  = PPTXColor(0x00, 0x79, 0x6B)
PPTX_WHITE = PPTXColor(0xFF, 0xFF, 0xFF)
PPTX_GRAY  = PPTXColor(0x61, 0x61, 0x61)
PPTX_LIGHT = PPTXColor(0xF0, 0xF4, 0xFF)


def _add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    from pptx.util import Emu
    from pptx.dml.color import RGBColor as RC
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def _text_box(slide, left, top, width, height, text, font_size=18,
               bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    from pptx.util import Pt
    tf = slide.shapes.add_textbox(left, top, width, height)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tf


def pptx_slide_title_only(prs, title_text, subtitle_text=None):
    """Add a slide with title + optional subtitle, blue background bar."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PPTX_LIGHT

    # Blue header bar
    _add_rect(slide, 0, 0, PPTX_W, Inches(1.5), PPTX_DARK)

    # Title in header bar
    _text_box(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(1.1),
              title_text, font_size=30, bold=True, color=PPTX_WHITE, align=PP_ALIGN.LEFT)

    if subtitle_text:
        _text_box(slide, Inches(0.4), Inches(1.6), Inches(12.5), Inches(0.6),
                  subtitle_text, font_size=14, bold=False, color=PPTX_GRAY, align=PP_ALIGN.LEFT)
    return slide


def pptx_add_bullets(slide, bullets, left=Inches(0.5), top=Inches(2.1),
                     width=Inches(12.3), height=Inches(4.5), font_size=14):
    from pptx.util import Pt
    tf_shape = slide.shapes.add_textbox(left, top, width, height)
    tf = tf_shape.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        # indent logic: sub-bullets start with spaces
        indent = b.startswith("  ")
        run = p.add_run()
        run.text = ("    • " if indent else "• ") + b.strip()
        run.font.size = Pt(font_size if not indent else font_size - 1)
        run.font.color.rgb = PPTX_GRAY if not indent else PPTXColor(0x42, 0x42, 0x42)


def pptx_add_two_col(slide, left_items, right_items, font_size=13):
    from pptx.util import Pt
    for items, left in [(left_items, Inches(0.5)), (right_items, Inches(6.8))]:
        tf_shape = slide.shapes.add_textbox(left, Inches(2.1), Inches(6.1), Inches(4.5))
        tf = tf_shape.text_frame
        tf.word_wrap = True
        first = True
        for b in items:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = "• " + b
            run.font.size = Pt(font_size)
            run.font.color.rgb = PPTX_GRAY


def pptx_add_step_boxes(slide, steps):
    """Render workflow steps as coloured boxes horizontally."""
    n = len(steps)
    box_w = Inches(12.0 / n)
    for i, (label, color) in enumerate(steps):
        left = Inches(0.5) + i * box_w
        _add_rect(slide, left, Inches(2.2), box_w - Inches(0.05), Inches(0.9), color)
        _text_box(slide, left, Inches(2.2), box_w - Inches(0.05), Inches(0.9),
                  label, font_size=11, bold=True, color=PPTX_WHITE, align=PP_ALIGN.CENTER)
        if i < n - 1:
            _text_box(slide, left + box_w - Inches(0.15), Inches(2.35), Inches(0.25), Inches(0.5),
                      "→", font_size=14, bold=True, color=PPTX_DARK, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# POWERPOINT DOCUMENT
# ─────────────────────────────────────────────────────────────────────────────

def build_pptx():
    from pptx import Presentation as Prs
    from pptx.util import Inches, Pt, Emu
    prs = Prs()
    prs.slide_width = PPTX_W
    prs.slide_height = PPTX_H

    # ── SLIDE 1: Title Slide ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = PPTX_DARK

    _text_box(sl, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.4),
              "Drug Management System", font_size=40, bold=True,
              color=PPTX_WHITE, align=PP_ALIGN.CENTER)
    _text_box(sl, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.8),
              "Adobe Acrobat Sign — User Guide by Role", font_size=22,
              bold=False, color=PPTXColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
    _text_box(sl, Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.5),
              f"Version {VERSION}  |  {TODAY}", font_size=13,
              color=PPTXColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

    # ── SLIDE 2: Overview of Adobe Sign in DMS ───────────────────────────────
    sl = pptx_slide_title_only(prs, "Adobe Acrobat Sign — System Overview")
    pptx_add_bullets(sl, [
        "Electronic signatures in DMS are powered by Adobe Acrobat Sign.",
        "Used for the final approval of drug management documents after the Approver reviews and accepts.",
        "NOT a manual wet-ink signature — Adobe Sign is a legally compliant electronic signature.",
        "Signature flow: Approver clicks 'Approve' → status moves to 'Pending Approval' → Adobe Sign is initiated.",
        "After signing in Adobe Sign, the Approver clicks 'Final Approve' → document becomes 'Signed / Final'.",
        "The signed PDF is automatically stored in the SharePoint Signed Documents library.",
        "",
        "Why Adobe Sign?",
        "  Legally binding under eIDAS (EU), ESIGN Act (USA), and other global regulations.",
        "  Full audit trail with timestamp, IP address, and signer identity.",
        "  Integrated directly with the DMS SharePoint environment.",
    ])

    # ── SLIDE 3: eSign by Role — Author ─────────────────────────────────────
    sl = pptx_slide_title_only(prs, "Adobe Sign — Author Role",
                                "What does the Author experience?")
    pptx_add_bullets(sl, [
        "The Author creates the document and submits it for review.",
        "The Author does NOT directly initiate the Adobe Sign process.",
        "After the Approver approves, the Author may receive an Adobe Sign notification email "
        "if the workflow is configured to require a countersignature.",
        "",
        "If countersignature is required:",
        "  Open the Adobe Sign email and click 'Review and Sign'.",
        "  The document opens in Adobe Sign. Review, then click to place your signature.",
        "  Click 'Click to Sign' to finalise.",
        "  The document status updates to 'Signed' once all parties have signed.",
        "",
        "If countersignature is NOT required:",
        "  The Author's involvement ends at submission. The Approver handles the signature.",
        "  The Author can track progress in 'My Documents' — the status badge updates automatically.",
    ])

    # ── SLIDE 4: eSign by Role — Approver ────────────────────────────────────
    sl = pptx_slide_title_only(prs, "Adobe Sign — Approver Role",
                                "What does the Approver experience?")
    pptx_add_bullets(sl, [
        "The Approver is the primary signer in the DMS Adobe Sign flow.",
        "",
        "Step-by-step experience:",
        "  1. Log in to SharePoint DMS.",
        "  2. Navigate to 'Assigned To Me' tab.",
        "  3. Open the document with status 'In Review'.",
        "  4. Review the document in the inline viewer. Add comments if needed.",
        "  5. Click 'Approve' — status changes to 'Pending Approval' and Adobe Sign is initiated.",
        "  6. Receive an Adobe Sign email with a direct link to the document.",
        "  7. Click the link — the Adobe Sign interface opens in your browser.",
        "  8. Review the document, then click the signature field. Type or draw your signature.",
        "  9. Click 'Click to Sign'. Return to DMS and click 'Final Approve'.",
        "  10. Document status updates to 'Signed / Final / Approved'.",
    ])

    # ── SLIDE 5: eSign by Role — Admin ───────────────────────────────────────
    sl = pptx_slide_title_only(prs, "Adobe Sign — Admin Role",
                                "What does the Admin experience?")
    pptx_add_bullets(sl, [
        "Admins do not typically participate in the signature flow.",
        "",
        "Admin capabilities:",
        "  View the status of any document across all stages in 'All Documents'.",
        "  Monitor signature status — look for 'Pending for Signature' or 'Signed' in the Status column.",
        "  Access the Signed Documents library in SharePoint to find all finalized PDF documents.",
        "  Manage users and ensure Approvers are correctly assigned to the DMS Approvers group.",
        "",
        "Where to find signed documents:",
        "  SharePoint Site → Signed Documents library (visible to Admin).",
        "  The signed PDF is automatically saved there once Adobe Sign confirms completion.",
        "  Document version is locked upon signing — no further edits are possible.",
    ])

    # ── SLIDE 6: How to Sign — Step by Step (Approver) ───────────────────────
    sl = pptx_slide_title_only(prs, "How to Sign — Step-by-Step (Approver)",
                                "Screen flow from login to completed signature")

    steps_top = [
        ("1. Log in to SharePoint DMS", PPTX_DARK),
        ("2. Open 'Assigned To Me'", PPTXColor(0x15, 0x65, 0xC0)),
        ("3. Click Document", PPTXColor(0x00, 0x79, 0x6B)),
        ("4. Click 'Approve'", PPTXColor(0xE6, 0x51, 0x00)),
    ]
    pptx_add_step_boxes(sl, steps_top)

    # Second row of steps
    steps_bot = [
        ("5. Adobe Sign Email", PPTX_DARK),
        ("6. Review & Sign", PPTXColor(0x15, 0x65, 0xC0)),
        ("7. Click to Sign ✓", PPTXColor(0x00, 0x79, 0x6B)),
        ("8. Click 'Final Approve'", PPTXColor(0x2E, 0x7D, 0x32)),
    ]
    n = len(steps_bot)
    box_w = Inches(12.0 / n)
    for i, (label, color) in enumerate(steps_bot):
        left = Inches(0.5) + i * box_w
        _add_rect(sl, left, Inches(3.4), box_w - Inches(0.05), Inches(0.9), color)
        _text_box(sl, left, Inches(3.4), box_w - Inches(0.05), Inches(0.9),
                  label, font_size=11, bold=True, color=PPTX_WHITE, align=PP_ALIGN.CENTER)

    pptx_add_bullets(sl, [
        "Status Transitions during this flow:",
        "  In Review → [Approve] → Pending Approval → [Adobe Sign] → Pending for Signature → [Final Approve] → Signed / Final",
    ], top=Inches(4.6), font_size=13)

    # ── SLIDE 7: Signed Documents Library ────────────────────────────────────
    sl = pptx_slide_title_only(prs, "Signed Documents Library",
                                "Where to find finalized signed documents")
    pptx_add_bullets(sl, [
        "Once a document is signed via Adobe Acrobat Sign:",
        "  The signed PDF is automatically retrieved and saved to the SharePoint 'Signed Documents' library.",
        "  The DMS document status is updated to 'Signed' or 'Final'.",
        "  The version is locked — no further edits can be made.",
        "",
        "How to access Signed Documents (Admin):",
        "  1. Navigate to the SharePoint site.",
        "  2. Open the 'Signed Documents' library from the site navigation or 'All Documents' → filter by 'Signed'.",
        "  3. Each signed PDF has metadata: Drug, Category, Approver, Signed Date.",
        "",
        "Authors and Approvers:",
        "  Authors can see the 'Signed' status badge in 'My Documents'.",
        "  Approvers can see the 'Signed' status in documents they approved.",
        "  Neither role can edit or delete the signed version.",
    ])

    # ── SLIDE 8: Summary ─────────────────────────────────────────────────────
    sl = pptx_slide_title_only(prs, "Summary — Adobe Sign in DMS",
                                "Key takeaways for all roles")
    pptx_add_two_col(sl,
        left_items=[
            "Adobe Acrobat Sign is the eSignature engine.",
            "Approver clicks 'Approve' to initiate — status: Pending Approval.",
            "Adobe Sign sends an email with a signing link.",
            "Approver signs, then clicks 'Final Approve' in DMS.",
            "Legally compliant signature with full audit trail.",
            "No additional software needed — works via email link.",
            "Author may countersign depending on configuration.",
            "Admin monitors status and accesses the Signed Documents library.",
        ],
        right_items=[
            "Document Status Flow:",
            "Draft → In Review → Pending Approval → Pending for Signature → Signed / Final",
            "",
            "Who signs:",
            "Approver (always) + Author (if countersignature configured).",
            "",
            "Where signed docs are stored:",
            "SharePoint 'Signed Documents' library — PDF, version-locked.",
        ],
        font_size=13
    )

    prs.save(PPTX_PATH)
    print(f"[OK] Adobe Sign PPT saved to: {PPTX_PATH}")


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating DMS SOP Operational Guide (Word)...")
    build_docx()
    print("Generating DMS Adobe Sign Guide (PowerPoint)...")
    build_pptx()
    print("\nDone. Files are in:", OUTPUT_DIR)
